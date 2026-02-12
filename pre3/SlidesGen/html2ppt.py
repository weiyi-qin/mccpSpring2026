#!/usr/bin/env python3
"""
Generate poeOutput.pptx with fully EDITABLE native PowerPoint shapes.

Recreates the UpANNS architecture diagram (poeOutput.html) as native shapes,
text boxes, connectors, and freeform polygons — all individually editable.

Usage:  python html2ppt.py
Requires: pip install python-pptx
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from lxml.etree import SubElement

SCRIPT_DIR = Path(__file__).resolve().parent
OUTPUT = SCRIPT_DIR / "poeOutput.pptx"

# ═══════════════════════════════════════════════════════════════
# Coordinate mapping:  SVG viewBox 1300×560 → Slide 13.333″×7.5″
# ═══════════════════════════════════════════════════════════════
S  = 0.01026     # uniform scale: inches per SVG unit
Y0 = 0.877       # vertical centering offset (inches)

def X(v):  return Inches(v * S)          # SVG x → slide x (EMU)
def Y(v):  return Inches(v * S + Y0)     # SVG y → slide y (EMU)
def W(v):  return Inches(v * S)          # SVG width → slide width
def H(v):  return Inches(v * S)          # SVG height → slide height

# ═══════════════════════════════════════════════════════════════
# Colour palette (matching SVG)
# ═══════════════════════════════════════════════════════════════
GREEN      = RGBColor(0x8f, 0xc4, 0x8a)
GOLD       = RGBColor(0xe8, 0xc9, 0x8a)
GREY       = RGBColor(0xbb, 0xbb, 0xbb)
LGREY      = RGBColor(0xdd, 0xdd, 0xdd)
DKGREEN    = RGBColor(0x2d, 0x6b, 0x2d)
ORANGE     = RGBColor(0xe8, 0xa0, 0x50)
WHITE      = RGBColor(0xff, 0xff, 0xff)
BLUE_BG    = RGBColor(0xd6, 0xe8, 0xf7)
BLUE_BRD   = RGBColor(0x7b, 0xad, 0xd4)
BLUE_TXT   = RGBColor(0x1a, 0x5a, 0x8a)
BLUE_DK    = RGBColor(0x4a, 0x90, 0xb8)
CPU_BG     = RGBColor(0xee, 0xee, 0xee)
CPU_BRD    = RGBColor(0xcc, 0xcc, 0xcc)
RED        = RGBColor(0xc0, 0x39, 0x2b)
DARK       = RGBColor(0x1a, 0x1a, 0x1a)
TXT        = RGBColor(0x33, 0x33, 0x33)
ARR        = RGBColor(0x44, 0x44, 0x44)
LN_CLR     = RGBColor(0x88, 0x88, 0x88)
STROKE     = RGBColor(0xaa, 0xaa, 0xaa)
CYL_GREEN  = RGBColor(0x6a, 0xb8, 0x6a)
CYL_ORANGE = RGBColor(0xe8, 0xa0, 0x50)
CYL_YELLOW = RGBColor(0xe8, 0xd8, 0x6a)
CYL_YGREEN = RGBColor(0xa8, 0xd8, 0x68)
DPU_FILL   = RGBColor(0x2a, 0x64, 0x96)
DPU_BRD    = RGBColor(0x1a, 0x4a, 0x72)
GREY66     = RGBColor(0x66, 0x66, 0x66)
GREY99     = RGBColor(0x99, 0x99, 0x99)
GREYBB     = RGBColor(0xbb, 0xbb, 0xbb)

# Grid colour look-up tables
CM  = {'g': GREEN, 'y': GOLD, 'x': GREY}
CM2 = {'d': LGREY, 'g': GREEN, 'y': GOLD, 'x': GREY}
CM3 = {'d': LGREY, 'g': GREEN, 'y': GOLD, 'dk': DKGREEN, 'o': ORANGE, 'x': GREY}

# ═══════════════════════════════════════════════════════════════
# Diagram data (extracted from SVG)
# ═══════════════════════════════════════════════════════════════

# Vector grids (8 cells each)
VECTORS = {
    0:   'g y x g x y g y'.split(),
    1:   'y x g x y g x y'.split(),
    2:   'x g y g x x y g'.split(),
    'n': 'g g x y g x y g'.split(),
}

# Co-occurrence grid (6 rows × 8 cols)
COOC = [
    'd d y d g d y d'.split(),
    'd g d x d d g y'.split(),
    'y d d g d y d x'.split(),
    'd d g d y d d g'.split(),
    'x d d y d x d d'.split(),
    'd y d d g d y g'.split(),
]

# Encoded blocks (right side of encoding box, 5 groups)
ENC_BLOCKS = [
    (50,  'd d y d'.split()),
    (85,  'g dk g d'.split()),
    (120, 'd o d'.split()),
    (155, 'x d'.split()),
    (185, ['y']),
]

# Merged result column
MERGED = 'd g dk o x y'.split()


# ═══════════════════════════════════════════════════════════════
# Helper functions
# ═══════════════════════════════════════════════════════════════

def _set_text(tf, text, font_size=Pt(10), font_bold=False,
              font_italic=False, font_color=DARK,
              font_name='Calibri', align=PP_ALIGN.CENTER):
    """Set (possibly multi-line) text on a TextFrame."""
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = font_bold
        run.font.italic = font_italic
        run.font.color.rgb = font_color
        run.font.name = font_name


def add_textbox(slide, sx, sy, sw, sh, text, **kw):
    """Add a text box at SVG coordinates."""
    box = slide.shapes.add_textbox(X(sx), Y(sy), W(sw), H(sh))
    _set_text(box.text_frame, text, **kw)
    return box


def add_rect(slide, sx, sy, sw, sh, fill=None, border=None,
             border_width=Pt(1), rx=0, text=None, **text_kw):
    """Add a rectangle (optionally rounded) with optional fill, border, text."""
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rx else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(st, X(sx), Y(sy), W(sw), H(sh))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    if text:
        _set_text(shape.text_frame, text, **text_kw)
        # Vertical-centre text inside shape
        body = shape._element.find(qn('p:txBody'))
        if body is not None:
            bp = body.find(qn('a:bodyPr'))
            if bp is not None:
                bp.set('anchor', 'ctr')
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ARR, width=Pt(1.2),
              end_arrow=True, start_arrow=False):
    """Add a connector line with optional arrowheads."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, X(x1), Y(y1), X(x2), Y(y2))
    conn.line.color.rgb = color
    conn.line.width = width
    # Arrowheads via OoXML
    spPr = conn._element.find(qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = SubElement(spPr, qn('a:ln'))
    if end_arrow:
        te = SubElement(ln, qn('a:tailEnd'))
        te.set('type', 'triangle')
        te.set('w', 'med')
        te.set('len', 'med')
    if start_arrow:
        he = SubElement(ln, qn('a:headEnd'))
        he.set('type', 'triangle')
        he.set('w', 'med')
        he.set('len', 'med')
    return conn


def add_line(slide, x1, y1, x2, y2, color=ARR, width=Pt(1),
             dashed=False):
    """Add a plain line (no arrowhead)."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, X(x1), Y(y1), X(x2), Y(y2))
    conn.line.color.rgb = color
    conn.line.width = width
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return conn


def add_sq(slide, sx, sy, size, fill, border=STROKE):
    """Add a tiny coloured square."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, X(sx), Y(sy), W(size), H(size))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(0.5)
    # remove default textbox content
    s.text_frame.clear()
    return s


def add_polygon(slide, points, fill=GREY, border=GREY99):
    """Add a closed freeform polygon (for thick directional arrows)."""
    ff = slide.shapes.build_freeform(X(points[0][0]), Y(points[0][1]))
    ff.add_line_segments([(X(p[0]), Y(p[1])) for p in points[1:]])
    s = ff.convert_to_shape()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(0.5)
    return s


def add_can(slide, sx, sy, sw, sh, fill):
    """Add a cylinder (CAN) shape."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.CAN, X(sx), Y(sy), W(sw), H(sh))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = RGBColor(0x55, 0x55, 0x55)
    s.line.width = Pt(0.7)
    s.text_frame.clear()
    return s


# ═══════════════════════════════════════════════════════════════
# Build the presentation
# ═══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ─────────────────────────────────────────────────────────
    # OFFLINE: Encoded Points + vector grids
    # ─────────────────────────────────────────────────────────
    add_textbox(slide, 5, 8, 140, 20, "Encoded points",
                font_size=Pt(11), font_bold=True, font_color=DARK,
                align=PP_ALIGN.CENTER)

    vec_rows = [(0, 36, 38), (1, 53, 55), (2, 70, 72)]
    for idx, grid_y, label_y in vec_rows:
        add_textbox(slide, 18, label_y, 55, 16, f"Vector {idx}",
                    font_size=Pt(9), font_italic=True, font_color=TXT,
                    align=PP_ALIGN.LEFT)
        for ci, c in enumerate(VECTORS[idx]):
            add_sq(slide, 75 + ci * 13, grid_y, 12, CM[c])

    add_textbox(slide, 88, 88, 24, 20, "\u22EE",
                font_size=Pt(14), font_bold=True, font_color=TXT,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, 18, 115, 55, 16, "Vector n",
                font_size=Pt(9), font_italic=True, font_color=TXT,
                align=PP_ALIGN.LEFT)
    for ci, c in enumerate(VECTORS['n']):
        add_sq(slide, 75 + ci * 13, 113, 12, CM[c])

    print("[1/8] Encoded points done")

    # ─────────────────────────────────────────────────────────
    # OFFLINE: Grey arrow → Co-occurrence Aware Encoding box
    # ─────────────────────────────────────────────────────────
    add_polygon(slide, [
        (182, 72), (222, 72), (222, 66), (236, 80),
        (222, 94), (222, 88), (182, 88)])

    # Outer box
    add_rect(slide, 245, 8, 380, 210,
             fill=BLUE_BG, border=BLUE_BRD, border_width=Pt(1.5), rx=6)
    add_textbox(slide, 310, 218, 250, 16,
                "Co-occurrence Aware Encoding",
                font_size=Pt(10), font_bold=True, font_color=DARK,
                align=PP_ALIGN.CENTER)

    # "High Frequent Co-occurrence" label
    add_rect(slide, 320, 16, 160, 32,
             fill=BLUE_DK, border=BLUE_DK, rx=4,
             text="High Frequent\nCo-occurrence",
             font_size=Pt(8), font_bold=True, font_color=WHITE)

    # Left grid (6×8, 14 px cells, 15 px step)
    for ri, row in enumerate(COOC):
        for ci, c in enumerate(row):
            add_sq(slide, 270 + ci * 15, 55 + ri * 15, 14, CM2[c], GREYBB)

    add_textbox(slide, 410, 118, 60, 16, "Encode",
                font_size=Pt(10), font_bold=True, font_color=DARK,
                align=PP_ALIGN.LEFT)

    # Arrows: grid → encoded blocks
    for x1, y1, x2, y2 in [
        (395, 100, 460, 65), (395, 120, 460, 100), (395, 140, 460, 135),
        (395, 160, 460, 165), (395, 175, 460, 195),
    ]:
        add_arrow(slide, x1, y1, x2, y2, LN_CLR, Pt(1))

    # Encoded blocks (right side, 5 groups)
    for gy, cols in ENC_BLOCKS:
        for ci, c in enumerate(cols):
            add_sq(slide, 465 + ci * 16, gy, 14, CM3[c], GREYBB)

    # Converging arrows → merged column
    for x1, y1, x2, y2 in [
        (535, 65, 570, 100), (535, 95, 570, 105), (530, 130, 570, 115),
        (520, 165, 570, 130), (510, 195, 570, 140),
    ]:
        add_arrow(slide, x1, y1, x2, y2, LN_CLR, Pt(1))

    # Merged result column (6 squares)
    for i, c in enumerate(MERGED):
        add_sq(slide, 575, 80 + i * 15, 14, CM3[c], GREYBB)

    print("[2/8] Co-occurrence encoding done")

    # ─────────────────────────────────────────────────────────
    # Grey arrow → Copy & Distribute
    # ─────────────────────────────────────────────────────────
    add_polygon(slide, [
        (630, 104), (670, 104), (670, 98), (684, 112),
        (670, 126), (670, 120), (630, 120)])

    # Outer box
    add_rect(slide, 695, 8, 580, 210,
             fill=BLUE_BG, border=BLUE_BRD, border_width=Pt(1.5), rx=6)
    add_textbox(slide, 885, 12, 200, 20, "Copy & Distribute",
                font_size=Pt(11), font_bold=True, font_color=BLUE_TXT,
                align=PP_ALIGN.CENTER)

    # Database cylinders (one tall CAN per column)
    for cx, color in [
        (720, CYL_GREEN), (775, CYL_ORANGE), (830, CYL_YELLOW),
        (920, CYL_GREEN), (975, CYL_YGREEN),
    ]:
        add_can(slide, cx, 42, 36, 66, color)

    add_textbox(slide, 878, 58, 32, 24, "\u00B7\u00B7\u00B7",
                font_size=Pt(16), font_bold=True, font_color=TXT,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, 890, 142, 190, 18, "Data Transfer",
                font_size=Pt(11), font_bold=True, font_color=BLUE_TXT,
                align=PP_ALIGN.CENTER)

    # Down arrows from cylinders to DPU chips
    for cx in [738, 793, 848, 938, 993]:
        add_arrow(slide, cx, 108, cx, 130, GREY66, Pt(1))

    # DPU chips (simplified as small blue rounded rects)
    dpu_xs = [720, 775, 830, 920, 975]
    for dx in dpu_xs:
        add_rect(slide, dx, 135, 36, 36,
                 fill=DPU_FILL, border=DPU_BRD, rx=2)

    # DPU labels
    for dx, label in zip(dpu_xs,
                         ["DPU 0", "DPU 1", "DPU 2", "DPU 0", "DPU n"]):
        add_textbox(slide, dx - 5, 175, 46, 14, label,
                    font_size=Pt(7), font_color=TXT,
                    align=PP_ALIGN.CENTER)

    add_textbox(slide, 873, 138, 32, 24, "\u00B7\u00B7\u00B7",
                font_size=Pt(12), font_bold=True, font_color=TXT,
                align=PP_ALIGN.CENTER)

    print("[3/8] Copy & Distribute done")

    # ─────────────────────────────────────────────────────────
    # DIVIDER  (Offline / Online)
    # ─────────────────────────────────────────────────────────
    add_line(slide, 0, 270, 1300, 270,
             GREY66, Pt(1.5), dashed=True)

    add_textbox(slide, 15, 248, 80, 20, "Offline",
                font_size=Pt(13), font_bold=True, font_italic=True,
                font_color=RED, align=PP_ALIGN.LEFT)
    add_textbox(slide, 15, 272, 80, 20, "Online",
                font_size=Pt(13), font_bold=True, font_italic=True,
                font_color=RED, align=PP_ALIGN.LEFT)

    print("[4/8] Divider done")

    # ─────────────────────────────────────────────────────────
    # ONLINE: CPU box + Query / Scheduling
    # ─────────────────────────────────────────────────────────
    add_rect(slide, 30, 305, 285, 195,
             fill=CPU_BG, border=CPU_BRD, border_width=Pt(1.5), rx=6)
    add_textbox(slide, 50, 306, 50, 18, "CPU",
                font_size=Pt(11), font_bold=True, font_color=TXT,
                align=PP_ALIGN.LEFT)

    # Query Batch
    add_rect(slide, 50, 430, 72, 42,
             fill=WHITE, border=GREY66, rx=4,
             text="Query\nBatch",
             font_size=Pt(8), font_bold=True, font_color=TXT)

    add_arrow(slide, 122, 451, 148, 451)

    # ①Cluster Filtering
    add_rect(slide, 152, 425, 100, 50,
             fill=BLUE_BG, border=BLUE_BRD, rx=4,
             text="\u2460Cluster\nFiltering",
             font_size=Pt(9), font_bold=True, font_color=BLUE_TXT)

    add_arrow(slide, 202, 425, 202, 385)

    # ②Query Scheduling
    add_rect(slide, 152, 335, 100, 48,
             fill=BLUE_BG, border=BLUE_BRD, rx=4,
             text="\u2461Query\nScheduling",
             font_size=Pt(9), font_bold=True, font_color=BLUE_TXT)

    add_arrow(slide, 252, 359, 310, 359)

    # Fan-out
    add_textbox(slide, 320, 386, 20, 24, "\u22EE",
                font_size=Pt(14), font_bold=True, font_color=TXT,
                align=PP_ALIGN.CENTER)
    add_arrow(slide, 310, 359, 340, 340)
    add_line(slide, 310, 359, 340, 359)
    add_arrow(slide, 335, 359, 375, 359)
    add_line(slide, 310, 359, 340, 420)
    add_arrow(slide, 340, 420, 340, 470)

    add_textbox(slide, 318, 475, 80, 16, "DPU \u00D7 n",
                font_size=Pt(9), font_bold=True, font_color=TXT,
                align=PP_ALIGN.LEFT)

    print("[5/8] CPU / Query section done")

    # ─────────────────────────────────────────────────────────
    # ONLINE: DPU Processing box + pipeline
    # ─────────────────────────────────────────────────────────
    add_rect(slide, 380, 300, 895, 220,
             fill=BLUE_BG, border=BLUE_BRD, border_width=Pt(1.5), rx=6)
    add_textbox(slide, 400, 302, 50, 18, "DPU",
                font_size=Pt(11), font_bold=True, font_color=BLUE_TXT,
                align=PP_ALIGN.LEFT)

    # Pipeline boxes
    pipeline = [
        (410, 330, 135, 52, "\u2462LUT\nConstruction"),
        (575, 330, 145, 52, "\u2463Distance\nCalculation"),
        (755, 330, 135, 52, "\u2464TOP-K\nPruning"),
        (930, 330, 145, 52, "\u2465Identifying\nTop-k"),
    ]
    for bx, by, bw, bh, label in pipeline:
        add_rect(slide, bx, by, bw, bh,
                 fill=WHITE, border=BLUE_BRD, rx=4,
                 text=label,
                 font_size=Pt(9), font_bold=True, font_color=BLUE_TXT)

    # Arrows between pipeline stages
    for x1, y1, x2, y2 in [
        (545, 356, 570, 356),
        (720, 356, 750, 356),
        (890, 356, 925, 356),
    ]:
        add_arrow(slide, x1, y1, x2, y2)

    # "Control" label
    add_textbox(slide, 748, 392, 70, 16, "Control",
                font_size=Pt(10), font_bold=True, font_italic=True,
                font_color=RED, align=PP_ALIGN.LEFT)

    # Double-headed arrows (pipeline ↔ resource management)
    add_arrow(slide, 647, 382, 647, 428, end_arrow=True, start_arrow=True)
    add_arrow(slide, 822, 382, 822, 428, end_arrow=True, start_arrow=True)

    # Efficient Resource Management box
    add_rect(slide, 490, 430, 475, 75,
             fill=BLUE_BG, border=BLUE_BRD, rx=5)
    add_textbox(slide, 550, 487, 360, 16,
                "Efficient Resource Management",
                font_size=Pt(10), font_bold=True, font_color=BLUE_TXT,
                align=PP_ALIGN.CENTER)

    # Sub-boxes
    add_rect(slide, 520, 440, 185, 38,
             fill=WHITE, border=BLUE_BRD, rx=4,
             text="Thread Scheduling",
             font_size=Pt(9), font_bold=True, font_color=BLUE_TXT)
    add_rect(slide, 740, 440, 195, 38,
             fill=WHITE, border=BLUE_BRD, rx=4,
             text="Memory Management",
             font_size=Pt(9), font_bold=True, font_color=BLUE_TXT)

    print("[6/8] DPU processing done")

    # ─────────────────────────────────────────────────────────
    # CAPTION
    # ─────────────────────────────────────────────────────────
    add_textbox(
        slide, 200, 535, 900, 22,
        "Figure 5: Overview of UpANNS. "
        "The grey blocks in each vector represent high frequency "
        "co-occurrences in the encoded vectors.",
        font_size=Pt(10), font_color=DARK, font_name='Times New Roman',
        align=PP_ALIGN.CENTER)

    print("[7/8] Caption done")

    # ─────────────────────────────────────────────────────────
    # Save
    # ─────────────────────────────────────────────────────────
    prs.save(str(OUTPUT))
    print(f"\n[8/8] Saved editable PowerPoint → {OUTPUT}")


if __name__ == "__main__":
    main()
